"""Vision helpers — multimodal extraction for PDF, PPTX, DOCX, and images.

**Status: placeholder for future release.**  The current v3 release uses
qwen2.5:7b (text-only).  A future version will add a vision-capable model
for image analysis, scanned-PDF OCR, and visual document understanding.

This module is the single place for all vision-related logic:
  - render_pdf_pages(): convert PDF pages to PIL Images
  - extract_pdf_images(): pull embedded images from a PDF
  - extract_pptx_images(): pull embedded images from a PPTX
  - extract_docx_images(): pull embedded images from a DOCX
  - is_weak_extraction(): heuristic to detect low-quality text extraction
  - image_to_base64() / file_to_base64(): encode images for Ollama's API
  - load_image_file(): load a standalone image file as a PIL Image
"""

from __future__ import annotations

import base64
import io
import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Supported standalone image extensions
# ---------------------------------------------------------------------------

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}

# ---------------------------------------------------------------------------
# Heuristic: detect "weak" text extraction from a PDF / PPTX / DOCX
# ---------------------------------------------------------------------------

_JUNK_RE = re.compile(r"[^\w\s]")
_MIN_CHARS_PER_PAGE = 40
_MAX_JUNK_RATIO = 0.70


def is_weak_extraction(text: str, page_count: int = 1) -> bool:
    """Return True when the extracted text is likely garbage or near-empty."""
    stripped = text.strip()
    if not stripped:
        return True
    if len(stripped) / max(page_count, 1) < _MIN_CHARS_PER_PAGE:
        return True
    junk_chars = len(_JUNK_RE.findall(stripped))
    if junk_chars / max(len(stripped), 1) > _MAX_JUNK_RATIO:
        return True
    return False


# ---------------------------------------------------------------------------
# PDF -> PIL Image rendering
# ---------------------------------------------------------------------------


def render_pdf_pages(path: Path, dpi: int = 150) -> list:
    """Render every page of a PDF to a list of PIL Images."""
    try:
        import fitz  # pymupdf
        from PIL import Image
    except ImportError:
        logger.debug("pymupdf or Pillow not installed — skipping PDF rendering")
        return []

    images: list = []
    try:
        doc = fitz.open(str(path))
        zoom = dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)
        for page in doc:
            pix = page.get_pixmap(matrix=matrix)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            images.append(img)
        doc.close()
    except Exception as exc:
        logger.warning("Failed to render PDF %s: %s", path, exc)
    return images


def extract_pdf_images(path: Path) -> list:
    """Extract embedded images from a PDF as PIL Images."""
    try:
        import fitz
        from PIL import Image
    except ImportError:
        return []

    images: list = []
    try:
        doc = fitz.open(str(path))
        for page in doc:
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                if base_image and base_image.get("image"):
                    img = Image.open(io.BytesIO(base_image["image"]))
                    if img.width >= 50 and img.height >= 50:
                        images.append(img)
        doc.close()
    except Exception as exc:
        logger.warning("Failed to extract images from PDF %s: %s", path, exc)
    return images


# ---------------------------------------------------------------------------
# PPTX image extraction
# ---------------------------------------------------------------------------


def extract_pptx_images(path: Path) -> list:
    """Extract embedded images from a PPTX as PIL Images."""
    try:
        from PIL import Image
        from pptx import Presentation
    except ImportError:
        return []

    images: list = []
    try:
        prs = Presentation(str(path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    blob = shape.image.blob
                    img = Image.open(io.BytesIO(blob))
                    if img.width >= 50 and img.height >= 50:
                        images.append(img)
    except Exception as exc:
        logger.warning("Failed to extract images from PPTX %s: %s", path, exc)
    return images


# ---------------------------------------------------------------------------
# DOCX image extraction
# ---------------------------------------------------------------------------


def extract_docx_images(path: Path) -> list:
    """Extract embedded images from a DOCX as PIL Images."""
    try:
        from docx import Document
        from PIL import Image
    except ImportError:
        return []

    images: list = []
    try:
        doc = Document(str(path))
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                blob = rel.target_part.blob
                img = Image.open(io.BytesIO(blob))
                if img.width >= 50 and img.height >= 50:
                    images.append(img)
    except Exception as exc:
        logger.warning("Failed to extract images from DOCX %s: %s", path, exc)
    return images


# ---------------------------------------------------------------------------
# Standalone image loading
# ---------------------------------------------------------------------------


def load_image_file(path: Path) -> "list":
    """Load a standalone image file and return it as a single-element list of PIL Images."""
    try:
        from PIL import Image
    except ImportError:
        return []
    try:
        img = Image.open(str(path))
        img.load()
        return [img]
    except Exception as exc:
        logger.warning("Failed to load image %s: %s", path, exc)
        return []


def is_image_file(path: Path) -> bool:
    """Return True if the path is a supported standalone image file."""
    return path.suffix.lower() in IMAGE_EXTENSIONS


# ---------------------------------------------------------------------------
# Image encoding for Ollama vision models
# ---------------------------------------------------------------------------


def image_to_base64(img, fmt: str = "PNG") -> str:
    """Encode a PIL Image to a base64 string suitable for Ollama's API."""
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    return base64.b64encode(buf.getvalue()).decode("utf-8")


def file_to_base64(path: Path) -> str:
    """Read a file from disk and return its base64 encoding."""
    return base64.b64encode(path.read_bytes()).decode("utf-8")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def pdf_page_count(path: Path) -> int:
    """Return the number of pages in a PDF (0 on error / missing dep)."""
    try:
        from pypdf import PdfReader
        return len(PdfReader(str(path)).pages)
    except Exception:
        return 0


def extract_images_from_file(path: Path) -> list:
    """Dispatch image extraction based on file extension."""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_images(path)
    if suffix == ".pptx":
        return extract_pptx_images(path)
    if suffix == ".docx":
        return extract_docx_images(path)
    if suffix in IMAGE_EXTENSIONS:
        return load_image_file(path)
    return []
